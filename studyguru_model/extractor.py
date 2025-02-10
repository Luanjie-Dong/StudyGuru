from pdf2image import convert_from_path
import pytesseract
from PyPDF2 import PdfReader
import fitz  
import re
from pptx import Presentation
from PIL import Image
from io import BytesIO
from pathlib import Path
import io

#StudyGuruExtractor = SGExtractor
class SGExtractor:

    def __init__(self,file_path):
        self.data = file_path
        


    def clean_text(self,content):
        content = re.sub(r"(\w+)-\n(\w+)", r"\1\2", content)
        unwanted_patterns = [
            "\\n",
            "  —",
            "——————————",
            "—————————",
            "—————",
            r"\\u[\dA-Fa-f]{4}",
            r"\uf075",
            r"\uf0b7",
        ]
        for pattern in unwanted_patterns:
            content = re.sub(pattern, "", content)

        content = re.sub(r"(\w)\s*-\s*(\w)", r"\1-\2", content)
        content = re.sub(r"\s+", " ", content)
        content = re.sub(r'[\u202a\u202b\u202c\u202d\u202e]', '', content)
        return content
    
    def read_pdf(self,pdf_path):
        doc = fitz.open(pdf_path)
        text_output = {}

        for page_num, page in enumerate(doc, start=1):
            print(f"Processing page {page_num}")

            text = self.clean_text(page.get_text("text"))
            if text:
                print("Extracted text using fitz.")
                text_output[page_num] = [text]

            images = page.get_images(full=True)
            images_text = ""
            for img in images:
                xref = img[0]  
                base_image = doc.extract_image(xref)  
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes)) 
                images_text += self.clean_text(pytesseract.image_to_string(image))
            text_output[page_num].append(images_text)

        return text_output
    
    def read_pptx(self, pptx_path):
        presentation = Presentation(pptx_path)
        text_output = {}

        for i, slide in enumerate(presentation.slides, start=1):
            print(f"Processing slide {i}")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    cleaned_text = self.clean_text(shape.text_frame.text) 
                    text_output.setdefault(i, []).append(cleaned_text) 

                if hasattr(shape, "image"):  
                    try:
                        image_bytes = shape.image.blob
                        image = Image.open(BytesIO(image_bytes))

                        ocr_text = self.clean_text(pytesseract.image_to_string(image))
                        text_output.setdefault(i, []).append(ocr_text)  
                    except Exception as e:
                        print(f"Error processing image on slide {i}: {e}")

        return text_output
    

    def read_image(self,image_path):
        image = Image.open(image_path)
        return {1:self.clean_text(pytesseract.image_to_string(image))}
    

    def extract_content(self):

        file_type = Path(self.data).suffix

        if file_type == ".pdf":
            return self.read_pdf(self.data)
        
        if file_type == ".pptx":
            return self.read_pptx(self.data)
        
        if file_type in (".jpg", ".png", ".bmp", ".tiff"):
            return self.read_image(self.data)
        
    
    def save_test_data(self,save_path):

        with open(save_path, "w", encoding="utf-8") as f:
            f.write(content)

        print(f"test data saved as {save_path}!")

        

        
    



if __name__=="__main__":

    test_type = "png"
    test_path = f'test_data/test.{test_type}'
    extractor = SGExtractor(test_path)
    content = extractor.extract_content()
    
    for page in content:
        print(page)
        print(content[page])
        print()

    # save_path = f"test_data/test_result({test_type}).txt"
    # save_test_content = extractor.save_test_data(save_path)

