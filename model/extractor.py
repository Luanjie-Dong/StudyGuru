from pdf2image import convert_from_path
import pytesseract
from PyPDF2 import PdfReader
import fitz  
import re
from pptx import Presentation
from PIL import Image
from io import BytesIO
from pathlib import Path
import io
import requests

#StudyGuruExtractor = SGExtractor
class SGExtractor:

    def __init__(self,file):

        self.file = file
        self.url = None
        self.data = None
        

    def get_content(self):

        try:
            api_host = "http://localhost:5005"
            url = f"{api_host}/one_note?note_id={self.file}"

            #Get Notes URL
            response = requests.get(url)
            response.raise_for_status()  
            json_data = response.json()
            self.url = json_data[0].get("pdf_URL")
            print(f"Successfully fetched url content for note ID: {self.file}")

            #Get Notes Content
            notes = requests.get(self.url)
            notes.raise_for_status()
            self.data = notes.content
            

        except requests.exceptions.RequestException as e:
            print(f"Error fetching content for note ID {self.file}: {e}")

    def clean_text(self,content):
        content = re.sub(r"(\w+)-\n(\w+)", r"\1\2", content)
        unwanted_patterns = [
            "\\n",
            "  —",
            "——————————",
            "—————————",
            "—————",
            r"\\u[\dA-Fa-f]{4}",
            r"\uf075",
            r"\uf0b7",
        ]
        for pattern in unwanted_patterns:
            content = re.sub(pattern, "", content)

        content = re.sub(r"(\w)\s*-\s*(\w)", r"\1-\2", content)
        content = re.sub(r"\s+", " ", content)
        content = re.sub(r'[\u202a\u202b\u202c\u202d\u202e]', '', content)
        return content
    
    def read_pdf(self):

        doc = fitz.open(stream=self.data, filetype="pdf")
        text_output = {}

        for page_num, page in enumerate(doc, start=1):
            print(f"Processing page {page_num}")

            text = self.clean_text(page.get_text("text"))
            if text:
                print("Extracted text using fitz.")
                text_output[page_num] = [text]

            images = page.get_images(full=True)
            images_text = ""
            for img in images:
                xref = img[0]  
                base_image = doc.extract_image(xref)  
                image_bytes = base_image["image"]
                image = Image.open(io.BytesIO(image_bytes)) 
                images_text += self.clean_text(pytesseract.image_to_string(image))
            text_output[page_num].append(images_text)

        return text_output
    
    def read_pptx(self):
        presentation = Presentation(BytesIO(self.data))
        text_output = {}

        for i, slide in enumerate(presentation.slides, start=1):
            print(f"Processing slide {i}")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    cleaned_text = self.clean_text(shape.text_frame.text) 
                    text_output.setdefault(i, []).append(cleaned_text) 

                if hasattr(shape, "image"):  
                    try:
                        image_bytes = shape.image.blob
                        image = Image.open(BytesIO(image_bytes))

                        ocr_text = self.clean_text(pytesseract.image_to_string(image))
                        text_output.setdefault(i, []).append(ocr_text)  
                    except Exception as e:
                        print(f"Error processing image on slide {i}: {e}")

        return text_output
    

    def read_image(self):
        image = Image.open(BytesIO(self.data))
        return {1:self.clean_text(pytesseract.image_to_string(image))}
    

    def extract_content(self):

        self.get_content()  

        file_type = Path(self.url).suffix.lower().strip("?")

        if file_type == ".pdf":
            return self.read_pdf()
        elif file_type == ".pptx":
            return self.read_pptx()
        elif file_type in (".jpg", ".png", ".bmp", ".tiff"):
            return self.read_image()
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
    
    def save_test_data(self,save_path):

        with open(save_path, "w", encoding="utf-8") as f:
            f.write(content)

        print(f"test data saved as {save_path}!")

        

        
    



if __name__=="__main__":

    test_type = "png"
    test_path = f'test_data/test.{test_type}'

    test_url = "https://lulvcodujqpxgvhkzyfc.supabase.co/storage/v1/object/public/notes/06ba3cc9-c8e5-4294-8e78-21cc6c7097d4/COGS_PwC_Case_Comp_Slidedeck.pdf?"
    extractor = SGExtractor(test_url)
    content = extractor.extract_content()
    
    for page in content:
        print(page)
        print(content[page])
        print()

    # save_path = f"test_data/test_result({test_type}).txt"
    # save_test_content = extractor.save_test_data(save_path)

